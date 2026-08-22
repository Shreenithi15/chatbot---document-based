import os
import re
import json
import numpy as np

# Suppress HuggingFace hub warnings and parallelism warnings
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["HF_HUB_DISABLE_SYMLINKS_WARNING"] = "1"

from pypdf import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
import faiss
from openai import OpenAI

# Optional dependencies for Word & PowerPoint documents
try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

# Initialize OpenAI client pointing to local LM Studio or API with timeout
client = OpenAI(
    base_url=os.getenv("OPENAI_BASE_URL", "http://127.0.0.1:1234/v1"),
    api_key=os.getenv("OPENAI_API_KEY", "lm-studio"),
    timeout=8.0
)

# Global embedding model cache
_EMBEDDING_MODEL = None

def get_embedding_model():
    global _EMBEDDING_MODEL
    if _EMBEDDING_MODEL is None:
        _EMBEDDING_MODEL = SentenceTransformer("all-MiniLM-L6-v2")
    return _EMBEDDING_MODEL

def read_pdf(pdf_path):
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"
    return text.strip()

def read_docx(docx_path):
    if docx is None:
        raise ImportError("python-docx is not installed.")
    doc = docx.Document(docx_path)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                full_text.append(" | ".join(row_text))
    return "\n".join(full_text)

def read_pptx(pptx_path):
    if pptx is None:
        raise ImportError("python-pptx is not installed.")
    prs = pptx.Presentation(pptx_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_runs.append(shape.text.strip())
    return "\n".join(text_runs)

def read_plain_text(file_path):
    encodings = ['utf-8', 'latin-1', 'cp1252', 'utf-16']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read().strip()
        except (UnicodeDecodeError, Exception):
            continue
    with open(file_path, 'rb') as f:
        return f.read().decode('utf-8', errors='ignore').strip()

def read_any_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        return read_pdf(file_path)
    elif ext in ['.docx', '.doc']:
        try:
            return read_docx(file_path)
        except Exception:
            return read_plain_text(file_path)
    elif ext in ['.pptx', '.ppt']:
        try:
            return read_pptx(file_path)
        except Exception:
            return read_plain_text(file_path)
    else:
        return read_plain_text(file_path)

def create_chunks(text):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )
    return splitter.split_text(text)

def create_faiss_index_from_chunks(chunks):
    model = get_embedding_model()
    embeddings = model.encode(chunks)
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(np.array(embeddings, dtype="float32"))
    return index, model

def extract_questions_fast_rules(user_input):
    """
    Ultra-fast rule-based question extractor (0.001 sec execution).
    """
    questions = []
    sentences = re.split(r'(?<=[.?!])\s+|\n+', user_input.strip())
    
    question_starters = (
        'what', 'how', 'why', 'who', 'where', 'when', 'which', 'whose', 'whom',
        'can', 'could', 'would', 'should', 'is', 'are', 'was', 'were', 'do', 'does', 'did',
        'explain', 'describe', 'list', 'detail', 'compare', 'tell', 'summarize'
    )

    for sentence in sentences:
        s = sentence.strip()
        if not s:
            continue
        s_lower = s.lower()
        if s.endswith('?') or any(s_lower.startswith(w) for w in question_starters):
            clean_q = s if s.endswith('?') else s + '?'
            if clean_q not in questions:
                questions.append(clean_q)

    return questions

def extract_questions(user_input):
    """
    Hybrid high-speed question extraction:
    First attempts ultra-fast rule parsing. If multiple questions are detected, returns instantly.
    Otherwise uses LLM for complex prompts.
    """
    if not user_input or not user_input.strip():
        return []

    # Fast path: rule-based check
    rule_qs = extract_questions_fast_rules(user_input)
    if len(rule_qs) >= 2:
        return rule_qs

    # LLM path for nuanced prompts
    try:
        response = client.chat.completions.create(
            model="qwen2.5-3b-instruct",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Extract ALL individual questions asked in the user text as a clean JSON array of strings. "
                        "Return ONLY valid JSON array format, e.g. [\"Question 1\", \"Question 2\"]."
                    )
                },
                {"role": "user", "content": user_input}
            ],
            temperature=0.1,
            max_tokens=150
        )
        content = response.choices[0].message.content.strip()
        if content.startswith("```"):
            content = re.sub(r"^```(?:json)?\n?", "", content)
            content = re.sub(r"\n?```$", "", content)
        
        parsed = json.loads(content)
        if isinstance(parsed, list) and len(parsed) > 0:
            return [str(q).strip() for q in parsed if str(q).strip()]
    except Exception:
        pass

    # Fallback to rule questions or full input
    if rule_qs:
        return rule_qs
    
    clean_input = user_input.strip()
    if not clean_input.endswith('?'):
        clean_input += '?'
    return [clean_input]

def extract_mentioned_document(user_input, available_filenames):
    """
    Scans user input for @filename or filename matches.
    Returns (cleaned_prompt, matched_filename).
    """
    if not available_filenames:
        return user_input, None

    # Check for explicit @mention pattern like @filename.pdf or @"filename.pdf" or @filename
    pattern = r'@(?:"([^"]+)"|([^\s@]+))'
    matches = re.findall(pattern, user_input)
    
    cleaned_input = user_input
    for match in matches:
        raw_target = match[0] or match[1]
        raw_target_clean = raw_target.strip().lower()
        
        for fname in available_filenames:
            fname_lower = fname.lower()
            fname_base = os.path.splitext(fname_lower)[0]
            if raw_target_clean == fname_lower or raw_target_clean == fname_base:
                # Remove mention tag from prompt text for cleaner QA
                cleaned_input = re.sub(rf'@(?:"{re.escape(raw_target)}"|{re.escape(raw_target)})', '', cleaned_input).strip()
                return cleaned_input, fname
                
    # Fallback: check if prompt mentions exact filename without @
    for fname in available_filenames:
        if fname.lower() in user_input.lower():
            return user_input, fname

    return user_input, None

def retrieve_context(question, model, index, chunks, k=2):
    if not index or not chunks:
        return ""
    question_embedding = model.encode([question])
    distances, indices = index.search(
        np.array(question_embedding, dtype="float32"),
        min(k, len(chunks))
    )
    retrieved_chunks = [chunks[i] for i in indices[0] if i < len(chunks)]
    return "\n\n".join(retrieved_chunks)

def retrieve_multi_context(question, model, documents_store, target_filename=None, k=2):
    """
    Retrieves relevant chunks across single or multiple stored documents.
    """
    if not documents_store or not model:
        return "", []

    docs_to_search = []
    if target_filename and target_filename in documents_store:
        docs_to_search = [target_filename]
    else:
        docs_to_search = list(documents_store.keys())

    all_retrieved = []  # list of (distance, chunk_text, filename)
    question_embedding = model.encode([question])
    q_vec = np.array(question_embedding, dtype="float32")

    for fname in docs_to_search:
        doc_data = documents_store[fname]
        index = doc_data.get("index")
        chunks = doc_data.get("chunks", [])
        if not index or not chunks:
            continue
        
        top_k = min(k, len(chunks))
        distances, indices = index.search(q_vec, top_k)
        for dist, idx in zip(distances[0], indices[0]):
            if idx < len(chunks):
                all_retrieved.append((float(dist), chunks[idx], fname))

    # Sort all retrieved chunks by L2 distance (ascending = best match first)
    all_retrieved.sort(key=lambda x: x[0])
    top_results = all_retrieved[:k * 2] if len(docs_to_search) > 1 else all_retrieved[:k]

    if not top_results:
        return "", []

    context_parts = []
    sources = set()
    for _, text, fname in top_results:
        context_parts.append(f"[{fname}]\n{text}")
        sources.add(fname)

    return "\n\n".join(context_parts), list(sources)

def answer_single_question(question, context=""):
    if context:
        prompt_content = f"Context:\n{context}\n\nQuestion:\n{question}"
        system_content = "Answer the question strictly and concisely using the provided context. Keep it direct and short."
    else:
        prompt_content = f"Question:\n{question}"
        system_content = "Answer the question concisely."

    try:
        response = client.chat.completions.create(
            model="qwen2.5-3b-instruct",
            messages=[
                {"role": "system", "content": system_content},
                {"role": "user", "content": prompt_content}
            ],
            temperature=0.3,
            max_tokens=250
        )
        return response.choices[0].message.content.strip()
    except Exception:
        if context:
            return f"Based on your document context:\n\n{context[:400]}..."
        else:
            return "Unable to connect to LLM server (LM Studio)."

def process_multi_question(user_input, chunks=None, index=None):
    extracted_qs = extract_questions(user_input)
    model = get_embedding_model() if (chunks and index) else None

    results = []
    for i, q in enumerate(extracted_qs, 1):
        context = ""
        if chunks and index and model:
            context = retrieve_context(q, model, index, chunks, k=2)

        answer = answer_single_question(q, context=context)
        results.append({
            "id": i,
            "question": q,
            "answer": answer,
            "context": context if context else None
        })

    return {
        "original_input": user_input,
        "question_count": len(results),
        "results": results
    }

def process_multi_question_store(user_input, documents_store, selected_doc=None):
    """
    Processes prompt using documents_store dict.
    Supports @mention targeting and fallback search.
    """
    available_files = list(documents_store.keys()) if documents_store else []
    
    target_doc = selected_doc if (selected_doc and selected_doc in documents_store) else None
    
    clean_input, mentioned_doc = extract_mentioned_document(user_input, available_files)
    if mentioned_doc:
        target_doc = mentioned_doc

    query_text = clean_input if clean_input.strip() else user_input
    extracted_qs = extract_questions(query_text)
    
    model = get_embedding_model() if documents_store else None

    results = []
    used_documents = set()

    for i, q in enumerate(extracted_qs, 1):
        context = ""
        sources = []
        if documents_store and model:
            context, sources = retrieve_multi_context(q, model, documents_store, target_filename=target_doc, k=2)
            used_documents.update(sources)

        answer = answer_single_question(q, context=context)
        results.append({
            "id": i,
            "question": q,
            "answer": answer,
            "context": context if context else None,
            "sources": sources
        })

    return {
        "original_input": user_input,
        "processed_input": query_text,
        "target_document": target_doc,
        "question_count": len(results),
        "results": results,
        "documents_used": list(used_documents)
    }

